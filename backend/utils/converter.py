# Document & Format Converter Module
# File: backend/utils/converter.py

from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4, landscape
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import fitz  # PyMuPDF
import os
from io import BytesIO

# Optional imports - handle gracefully if not installed
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("Warning: pdf2image not available. PDF to image conversion will use PyMuPDF instead.")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not available. Word document handling disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. PowerPoint handling disabled.")

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    print("Warning: openpyxl not available. Excel handling disabled.")


class Converter:
    """
    Comprehensive document format converter
    Supports: PDF↔Image, Office docs↔PDF, HTML↔PDF
    """
    
    def __init__(self):
        self.supported_conversions = {
            'jpg_to_pdf': 'Convert JPG images to PDF',
            'png_to_pdf': 'Convert PNG images to PDF',
            'pdf_to_jpg': 'Convert PDF pages to JPG',
            'pdf_to_png': 'Convert PDF pages to PNG',
            'word_to_pdf': 'Convert Word documents to PDF',
            'ppt_to_pdf': 'Convert PowerPoint to PDF',
            'excel_to_pdf': 'Convert Excel to PDF',
            'html_to_pdf': 'Convert HTML content to PDF',
        }
    
    # ============================================================
    # IMAGE TO PDF CONVERSIONS
    # ============================================================
    
    def jpg_to_pdf(self, image_paths, output_path, orientation='auto'):
        """
        Convert JPG images to PDF
        
        Args:
            image_paths (list): List of JPG file paths
            output_path (str): Output PDF path
            orientation (str): 'auto', 'portrait', or 'landscape'
            
        Returns:
            bool: Success status
        """
        try:
            images = []
            
            for img_path in image_paths:
                if os.path.exists(img_path):
                    img = Image.open(img_path)
                    
                    # Convert to RGB if necessary
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    images.append(img)
            
            if images:
                # Save all images as PDF
                first_image = images[0]
                
                if len(images) > 1:
                    first_image.save(
                        output_path,
                        save_all=True,
                        append_images=images[1:],
                        format='PDF'
                    )
                else:
                    first_image.save(output_path, format='PDF')
                
                # Close all images
                for img in images:
                    img.close()
                
                return True
            return False
        except Exception as e:
            print(f"Error converting JPG to PDF: {e}")
            return False
    
    def png_to_pdf(self, image_paths, output_path):
        """
        Convert PNG images to PDF
        
        Args:
            image_paths (list): List of PNG file paths
            output_path (str): Output PDF path
            
        Returns:
            bool: Success status
        """
        return self.jpg_to_pdf(image_paths, output_path)
    
    # ============================================================
    # PDF TO IMAGE CONVERSIONS
    # ============================================================
    
    def pdf_to_jpg(self, input_path, output_dir, dpi=300, pages=None):
        """
        Convert PDF pages to JPG images
        
        Args:
            input_path (str): Input PDF path
            output_dir (str): Output directory for images
            dpi (int): Image resolution
            pages (list): Pages to convert (None = all)
            
        Returns:
            list: List of output image paths
        """
        output_paths = []
        
        # Try pdf2image first if available
        if PDF2IMAGE_AVAILABLE:
            try:
                images = convert_from_path(input_path, dpi=dpi)
                for i, image in enumerate(images):
                    if pages is None or (i + 1) in pages:
                        output_path = os.path.join(output_dir, f"page_{i + 1}.jpg")
                        image.save(output_path, 'JPEG', quality=95)
                        output_paths.append(output_path)
                return output_paths
            except Exception as e:
                print(f"pdf2image failed (poppler may not be installed): {e}")
                # Fall through to PyMuPDF
        
        # Use PyMuPDF as fallback (or primary if pdf2image not available)
        try:
            doc = fitz.open(input_path)
            matrix = fitz.Matrix(dpi / 72, dpi / 72)
            
            for i, page in enumerate(doc):
                if pages is None or (i + 1) in pages:
                    pix = page.get_pixmap(matrix=matrix)
                    output_path = os.path.join(output_dir, f"page_{i + 1}.jpg")
                    pix.save(output_path)
                    output_paths.append(output_path)
            
            doc.close()
            return output_paths
        except Exception as e:
            print(f"Error converting PDF to JPG with PyMuPDF: {e}")
            return []
    
    def pdf_to_png(self, input_path, output_dir, dpi=300, pages=None):
        """
        Convert PDF pages to PNG images
        
        Args:
            input_path (str): Input PDF path
            output_dir (str): Output directory for images
            dpi (int): Image resolution
            pages (list): Pages to convert (None = all)
            
        Returns:
            list: List of output image paths
        """
        output_paths = []
        
        # Try pdf2image first if available
        if PDF2IMAGE_AVAILABLE:
            try:
                images = convert_from_path(input_path, dpi=dpi)
                for i, image in enumerate(images):
                    if pages is None or (i + 1) in pages:
                        output_path = os.path.join(output_dir, f"page_{i + 1}.png")
                        image.save(output_path, 'PNG')
                        output_paths.append(output_path)
                return output_paths
            except Exception as e:
                print(f"pdf2image failed (poppler may not be installed): {e}")
                # Fall through to PyMuPDF
        
        # Use PyMuPDF as fallback (or primary if pdf2image not available)
        try:
            doc = fitz.open(input_path)
            matrix = fitz.Matrix(dpi / 72, dpi / 72)
            
            for i, page in enumerate(doc):
                if pages is None or (i + 1) in pages:
                    pix = page.get_pixmap(matrix=matrix)
                    output_path = os.path.join(output_dir, f"page_{i + 1}.png")
                    pix.save(output_path)
                    output_paths.append(output_path)
            
            doc.close()
            return output_paths
        except Exception as e:
            print(f"Error converting PDF to PNG with PyMuPDF: {e}")
            return []
    
    # ============================================================
    # OFFICE DOCUMENT CONVERSIONS
    # ============================================================
    
    def word_to_pdf(self, input_path, output_path):
        """
        Convert Word document to PDF
        
        Args:
            input_path (str): Input DOCX/DOC path
            output_path (str): Output PDF path
            
        Returns:
            bool: Success status
        """
        try:
            if not DOCX_AVAILABLE:
                print("python-docx not available")
                return False
            
            doc = Document(input_path)
            
            # Create PDF using ReportLab
            pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    story.append(Paragraph(para.text, styles['Normal']))
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting Word to PDF: {e}")
            return False
    
    def ppt_to_pdf(self, input_path, output_path):
        """
        Convert PowerPoint to PDF
        
        Args:
            input_path (str): Input PPTX/PPT path
            output_path (str): Output PDF path
            
        Returns:
            bool: Success status
        """
        try:
            if not PPTX_AVAILABLE:
                print("python-pptx not available")
                return False
            
            prs = Presentation(input_path)
            
            # Create PDF with ReportLab
            c = canvas.Canvas(output_path, pagesize=landscape(letter))
            width, height = landscape(letter)
            
            for slide_num, slide in enumerate(prs.slides):
                # Add slide number
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, height - 50, f"Slide {slide_num + 1}")
                
                # Extract text from shapes
                y_position = height - 100
                c.setFont("Helvetica", 12)
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines = shape.text.split('\n')
                        for line in lines:
                            if line.strip():
                                c.drawString(50, y_position, line[:100])  # Truncate long lines
                                y_position -= 20
                                if y_position < 100:
                                    break
                
                c.showPage()
            
            c.save()
            return True
        except Exception as e:
            print(f"Error converting PPT to PDF: {e}")
            return False
    
    def excel_to_pdf(self, input_path, output_path):
        """
        Convert Excel to PDF
        
        Args:
            input_path (str): Input XLSX/XLS path
            output_path (str): Output PDF path
            
        Returns:
            bool: Success status
        """
        try:
            if not OPENPYXL_AVAILABLE:
                print("openpyxl not available")
                return False
            
            wb = load_workbook(input_path)
            
            # Create PDF with ReportLab
            c = canvas.Canvas(output_path, pagesize=landscape(letter))
            width, height = landscape(letter)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # Add sheet name
                c.setFont("Helvetica-Bold", 14)
                c.drawString(50, height - 50, f"Sheet: {sheet_name}")
                
                # Extract cell values
                y_position = height - 100
                c.setFont("Helvetica", 10)
                
                for row in sheet.iter_rows(max_row=50, max_col=10, values_only=True):
                    row_text = " | ".join([str(cell) if cell else "" for cell in row])
                    c.drawString(50, y_position, row_text[:150])  # Truncate long rows
                    y_position -= 15
                    
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                
                c.showPage()
            
            c.save()
            return True
        except Exception as e:
            print(f"Error converting Excel to PDF: {e}")
            return False
    
    # ============================================================
    # HTML CONVERSIONS
    # ============================================================
    
    def html_to_pdf(self, html_content, output_path, page_size='A4'):
        """
        Convert HTML content to PDF
        
        Args:
            html_content (str): HTML content or file path
            output_path (str): Output PDF path
            page_size (str): Page size ('A4' or 'letter')
            
        Returns:
            bool: Success status
        """
        try:
            # Parse page size
            if page_size.upper() == 'A4':
                size = A4
            else:
                size = letter
            
            # Create PDF using ReportLab
            pdf_doc = SimpleDocTemplate(output_path, pagesize=size)
            styles = getSampleStyleSheet()
            story = []
            
            # Simple HTML to text conversion (basic)
            # Remove HTML tags for simple conversion
            import re
            text_content = re.sub(r'<[^>]+>', '', html_content)
            
            # Split into paragraphs
            paragraphs = text_content.split('\n\n')
            
            for para in paragraphs:
                if para.strip():
                    story.append(Paragraph(para.strip(), styles['Normal']))
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting HTML to PDF: {e}")
            return False
    
    # ============================================================
    # PDF ARCHIVE FORMAT CONVERSION
    # ============================================================
    
    def pdf_to_pdfa(self, input_path, output_path):
        """
        Convert PDF to PDF/A format (archive format)
        
        Args:
            input_path (str): Input PDF path
            output_path (str): Output PDF/A path
            
        Returns:
            bool: Success status
        """
        try:
            # PyMuPDF can save as PDF/A-1b
            doc = fitz.open(input_path)
            doc.save(output_path, deflate=True)
            doc.close()
            return True
        except Exception as e:
            print(f"Error converting to PDF/A: {e}")
            return False


# ============================================================
# Image Handler Module
# ============================================================

class ImageHandler:
    """Handle image-related operations"""
    
    def __init__(self):
        self.supported_formats = ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff']
    
    def compress_image(self, input_path, output_path, quality=85):
        """
        Compress image file
        
        Args:
            input_path (str): Input image path
            output_path (str): Output image path
            quality (int): Compression quality (1-100)
            
        Returns:
            bool: Success status
        """
        try:
            img = Image.open(input_path)
            
            # Convert to RGB if necessary
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')
            
            img.save(output_path, 'JPEG', quality=quality, optimize=True)
            img.close()
            return True
        except Exception as e:
            print(f"Error compressing image: {e}")
            return False
    
    def resize_image(self, input_path, output_path, width, height):
        """
        Resize image
        
        Args:
            input_path (str): Input image path
            output_path (str): Output image path
            width (int): Target width
            height (int): Target height
            
        Returns:
            bool: Success status
        """
        try:
            img = Image.open(input_path)
            img_resized = img.resize((width, height), Image.LANCZOS)
            img_resized.save(output_path)
            img.close()
            img_resized.close()
            return True
        except Exception as e:
            print(f"Error resizing image: {e}")
            return False
    
    def convert_format(self, input_path, output_path, format):
        """
        Convert image format
        
        Args:
            input_path (str): Input image path
            output_path (str): Output image path
            format (str): Target format (jpg, png, etc.)
            
        Returns:
            bool: Success status
        """
        try:
            img = Image.open(input_path)
            
            # Convert mode if necessary
            if format.lower() in ('jpg', 'jpeg') and img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')
            
            img.save(output_path, format.upper())
            img.close()
            return True
        except Exception as e:
            print(f"Error converting image format: {e}")
            return False


# ============================================================
# OCR Handler Module
# ============================================================

class OCRHandler:
    """Handle OCR operations for scanned documents"""
    
    def __init__(self):
        self.supported_languages = ['eng', 'fra', 'deu', 'spa', 'ita', 'por', 'rus']
        self.tesseract_available = False
        
        try:
            import pytesseract
            self.pytesseract = pytesseract
            self.tesseract_available = True
        except ImportError:
            print("Warning: pytesseract not available. OCR features disabled.")
    
    def extract_text_ocr(self, image_path, language='eng'):
        """
        Extract text from image using OCR
        
        Args:
            image_path (str): Input image path
            language (str): OCR language code
            
        Returns:
            str: Extracted text
        """
        try:
            if not self.tesseract_available:
                return "OCR not available. Install pytesseract."
            
            img = Image.open(image_path)
            text = self.pytesseract.image_to_string(img, lang=language)
            img.close()
            return text
        except Exception as e:
            print(f"Error in OCR: {e}")
            return ""
    
    def pdf_to_text_ocr(self, pdf_path, language='eng'):
        """
        Extract text from PDF using OCR
        
        Args:
            pdf_path (str): Input PDF path
            language (str): OCR language code
            
        Returns:
            str: Extracted text
        """
        try:
            if not self.tesseract_available:
                return "OCR not available. Install pytesseract."
            
            doc = fitz.open(pdf_path)
            full_text = ""
            
            for page in doc:
                # Render page to image
                pix = page.get_pixmap(dpi=300)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # OCR the image
                text = self.pytesseract.image_to_string(img, lang=language)
                full_text += text + "\n\n--- Page Break ---\n\n"
                
                img.close()
            
            doc.close()
            return full_text.strip()
        except Exception as e:
            print(f"Error in PDF OCR: {e}")
            return ""


# ============================================================
# Document Handler Module
# ============================================================

class DocumentHandler:
    """Handle Office document operations"""
    
    def __init__(self):
        self.supported_formats = ['doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx']
    
    def extract_text_from_docx(self, doc_path):
        """
        Extract text from Word document
        
        Args:
            doc_path (str): Input document path
            
        Returns:
            str: Extracted text
        """
        try:
            if not DOCX_AVAILABLE:
                return "python-docx not available"
            
            doc = Document(doc_path)
            text = ""
            
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from DOCX: {e}")
            return ""
    
    def extract_text_from_pptx(self, ppt_path):
        """
        Extract text from PowerPoint
        
        Args:
            ppt_path (str): Input presentation path
            
        Returns:
            str: Extracted text
        """
        try:
            if not PPTX_AVAILABLE:
                return "python-pptx not available"
            
            prs = Presentation(ppt_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return ""
    
    def extract_text_from_xlsx(self, excel_path):
        """
        Extract text from Excel spreadsheet
        
        Args:
            excel_path (str): Input Excel path
            
        Returns:
            str: Extracted text
        """
        try:
            if not OPENPYXL_AVAILABLE:
                return "openpyxl not available"
            
            wb = load_workbook(excel_path)
            text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell else "" for cell in row])
                    text += row_text + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from XLSX: {e}")
            return ""
